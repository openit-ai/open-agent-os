"""Attachment text extraction — lazy deps, no hard failures.

Handles:
  pdf  -> pypdf (PdfReader) fallback pdfminer.six
  docx -> python-docx
  xlsx -> openpyxl
  pptx -> python-pptx
  txt/md/csv/json -> plain read
  images (png/jpg/jpeg/gif/webp/bmp/tiff) -> attachment reference + runtime user-turn instruction (NO OCR, NO LLM selection)

All imports are inside functions (lazy). If lib missing, fallback to "" or stub string.

Image handling (corrected per user instruction):
  - Do NOT select any LLM / model / provider.
  - Do NOT call any separate OCR or vision API nor standalone vision endpoint.
  - Image is represented as an attachment reference dict and a user-turn
    instruction string to be forwarded through the currently active Agent
    Runtime conversation context (session_id/tenant_id/user_id/channel/root/post)
    via the existing OAOS ACP / Hermes path (control_plane.acp_adapter.ACPAdapter).
  - No provider / model / LLM endpoint is invented in this package.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

SUPPORTED_EXTENSIONS = frozenset({
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    ".txt", ".md", ".csv", ".json",
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif",
})

_IMAGE_EXTS = frozenset({".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"})
_MAX_CHARS_DEFAULT = 20000

# Marker used in runtime instruction so tests can assert it's a runtime prompt, not OCR text
IMAGE_RUNTIME_INSTRUCTION_MARKER = "Image attachment reference"
IMAGE_RUNTIME_INSTRUCTION_KIND = "image_attachment_reference"

_MIME_BY_EXT = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
    ".tiff": "image/tiff",
    ".tif": "image/tiff",
}


def _truncate(text: str, max_chars: int) -> str:
    if max_chars and len(text) > max_chars:
        return text[:max_chars] + f"\n...[truncated {len(text)-max_chars} chars]"
    return text


def _read_txt(path: Path, max_chars: int) -> str:
    try:
        try:
            t = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            t = path.read_text(encoding="latin-1")
        return _truncate(t, max_chars)
    except Exception as e:
        return f"[txt read error: {e}]"


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def extract_pdf(path: Path | str, max_chars: int = _MAX_CHARS_DEFAULT) -> str:
    p = Path(path)
    try:
        from pypdf import PdfReader  # type: ignore
        reader = PdfReader(str(p))
        parts: list[str] = []
        for page in getattr(reader, "pages", []):
            try:
                txt = page.extract_text() or ""
                if txt:
                    parts.append(txt)
            except Exception:
                continue
        if parts:
            return _truncate("\n".join(parts), max_chars)
    except Exception:
        pass
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract  # type: ignore
        txt = pdfminer_extract(str(p)) or ""
        if txt.strip():
            return _truncate(txt, max_chars)
    except Exception:
        pass
    try:
        size = p.stat().st_size
        return f"[pdf extraction unavailable — install pypdf or pdfminer.six (file {p.name}, {size} bytes)]"
    except Exception as e:
        return f"[pdf extraction failed: {e}]"


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def extract_docx(path: Path | str, max_chars: int = _MAX_CHARS_DEFAULT) -> str:
    p = Path(path)
    try:
        from docx import Document  # type: ignore
        doc = Document(str(p))
        parts: list[str] = []
        for para in getattr(doc, "paragraphs", []):
            t = getattr(para, "text", "") or ""
            if t:
                parts.append(t)
        for table in getattr(doc, "tables", []):
            for row in getattr(table, "rows", []):
                cells = [getattr(c, "text", "") for c in getattr(row, "cells", [])]
                if any(cells):
                    parts.append(" | ".join(cells))
        if parts:
            return _truncate("\n".join(parts), max_chars)
        return ""
    except Exception as e:
        try:
            import importlib.util
            if importlib.util.find_spec("docx") is None:
                return f"[docx extraction unavailable — install python-docx (file {p.name})]"
        except Exception:
            pass
        return f"[docx extraction failed: {e}]"


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

def extract_xlsx(path: Path | str, max_chars: int = _MAX_CHARS_DEFAULT) -> str:
    p = Path(path)
    try:
        from openpyxl import load_workbook  # type: ignore
        wb = load_workbook(str(p), read_only=True, data_only=True)
        out: list[str] = []
        for ws in wb.worksheets:
            out.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) if v is not None else "" for v in row]
                if any(v.strip() for v in vals):
                    out.append("\t".join(vals))
            out.append("")
        wb.close()
        text = "\n".join(out).strip()
        return _truncate(text, max_chars) if text else ""
    except Exception as e:
        try:
            import importlib.util
            if importlib.util.find_spec("openpyxl") is None:
                return f"[xlsx extraction unavailable — install openpyxl (file {p.name})]"
        except Exception:
            pass
        return f"[xlsx extraction failed: {e}]"


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def extract_pptx(path: Path | str, max_chars: int = _MAX_CHARS_DEFAULT) -> str:
    p = Path(path)
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(p))
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, 1):  # type: ignore[attr-defined]
            parts.append(f"# Slide {idx}")
            for shape in slide.shapes:  # type: ignore[attr-defined]
                try:
                    txt = getattr(shape, "text", "") or ""
                    if txt:
                        parts.append(txt)
                except Exception:
                    pass
                if getattr(shape, "has_table", False):  # type: ignore[attr-defined]
                    try:
                        for row in shape.table.rows:  # type: ignore
                            cells = [getattr(c, "text", "") for c in row.cells]  # type: ignore
                            if any(cells):
                                parts.append(" | ".join(cells))
                    except Exception:
                        pass
            parts.append("")
        text = "\n".join(parts).strip()
        return _truncate(text, max_chars) if text else ""
    except Exception as e:
        try:
            import importlib.util
            if importlib.util.find_spec("pptx") is None:
                return f"[pptx extraction unavailable — install python-pptx (file {p.name})]"
        except Exception:
            pass
        return f"[pptx extraction failed: {e}]"


# ---------------------------------------------------------------------------
# Images — attachment reference + runtime instruction (NO OCR, NO LLM selection)
# ---------------------------------------------------------------------------

def is_image(path: Path | str) -> bool:
    """Return True if path has an image extension (no FS check)."""
    return Path(path).suffix.lower() in _IMAGE_EXTS


def build_image_attachment_reference(
    path: Path | str,
    vault_path: str | None = None,
    attachment_id: str | None = None,
) -> dict[str, Any]:
    """Build attachment reference dict for an image (no OCR, no model selection).

    This dict is intended to be forwarded as a user-turn attachment via
    the active Agent Runtime (ACP/Hermes) — not processed locally.
    """
    p = Path(path)
    ext = p.suffix.lower()
    mime = _MIME_BY_EXT.get(ext, "application/octet-stream")
    try:
        size = p.stat().st_size if p.exists() else 0
    except Exception:
        size = 0
    # vault_path must be tenant/agent isolated; never expose absolute filesystem path
    # fallback to sanitized filename only (not absolute str(p)). Reject an
    # absolute caller-provided path as it would break the ownership boundary.
    if vault_path and Path(vault_path).is_absolute():
        raise ValueError("vault_path must be relative to the owner-scoped vault")
    safe_vault = vault_path if vault_path else p.name
    ref: dict[str, Any] = {
        "kind": IMAGE_RUNTIME_INSTRUCTION_KIND,
        "filename": p.name,
        "ext": ext,
        "mime": mime,
        "size_bytes": size,
        "vault_path": safe_vault,
    }
    if attachment_id:
        ref["attachment_id"] = attachment_id
    return ref


def build_image_runtime_instruction(
    path: Path | str,
    vault_path: str | None = None,
    attachment_id: str | None = None,
    extra_instruction: str | None = None,
) -> str:
    """Build a user-turn instruction to be sent through active Agent Runtime.

    The returned string is a prompt for the currently active Hermes/ACP
    runtime — NOT a stored OCR result. No provider/model/LLM is selected here.
    """
    p = Path(path)
    ext = p.suffix.lower()
    try:
        size = p.stat().st_size if p.exists() else 0
    except Exception:
        size = 0
    # never expose absolute path; use isolated vault_path or plain filename
    vp = vault_path if vault_path else p.name
    # JSON-safe metadata encoding (prevents breaking on ', \", newline, etc.)
    metadata_json = json.dumps({"filename": p.name, "ext": ext, "bytes": size}, ensure_ascii=False)
    base = (
        f"[Image Attachment \u2014 LLM Vision Pending] {p.name} ({ext}, {size} bytes) stored at {vp} — attachment_id={attachment_id or 'pending'}\n"
        f"[{IMAGE_RUNTIME_INSTRUCTION_MARKER}: {p.name} ({ext}, {size} bytes) stored at {vp}]\n"
        "LLM Vision Request Prompt (deterministic, no local OCR):\n"
        "Please analyze this image attachment via the LLM vision model at query/runtime time. "
        "Describe, transcribe, or interpret the image as requested by the user. "
        "Status: pending LLM vision processing \u2014 no OCR text extracted locally; no separate OCR or vision API was invoked.\n"
        f"Metadata: {metadata_json}\n"
        "This is a user-turn instruction forwarded via the active Agent Runtime (ACP/Hermes)."
    )
    if extra_instruction:
        base = base + " " + extra_instruction.strip()
    return base


def extract_image(path: Path | str, max_chars: int = _MAX_CHARS_DEFAULT) -> str:
    """Image handler: return attachment-reference runtime instruction (no OCR).

    Explicitly does NOT import/call pytesseract, easyocr, ollama, or any
    vision API. No LLM/model/provider is selected.
    """
    # No lazy OCR imports here — intentionally absent.
    return _truncate(build_image_runtime_instruction(path), max_chars)


# ---------------------------------------------------------------------------
# Main dispatcher
# ---------------------------------------------------------------------------

def extract_text(path: Path | str, max_chars: int = _MAX_CHARS_DEFAULT) -> str:
    """Extract text from attachment by extension (lazy deps, no hard failure).

    For images: returns attachment-reference runtime instruction (no OCR).
    No provider/model selection is performed.
    """
    p = Path(path)
    ext = p.suffix.lower()
    if not p.exists():
        return f"[file not found: {p}]"
    if ext == ".pdf":
        return extract_pdf(p, max_chars)
    if ext in (".docx", ".doc"):
        return extract_docx(p, max_chars)
    if ext in (".xlsx", ".xls"):
        return extract_xlsx(p, max_chars)
    if ext in (".pptx", ".ppt"):
        return extract_pptx(p, max_chars)
    if ext in _IMAGE_EXTS:
        return extract_image(p, max_chars)
    if ext in (".txt", ".md", ".csv", ".json"):
        return _read_txt(p, max_chars)
    try:
        txt = _read_txt(p, max_chars)
        if txt and len(txt) < max_chars:
            return txt
        if ext not in SUPPORTED_EXTENSIONS:
            return f"[unsupported extension {ext} — treated as binary, size {p.stat().st_size} bytes]"
        return txt
    except Exception as e:
        return f"[extraction failed for {p.name}: {e}]"


# Alias
extract_attachment = extract_text


def is_supported(path: Path | str) -> bool:
    return Path(path).suffix.lower() in SUPPORTED_EXTENSIONS
